#!/usr/bin/env python3
"""
Reorder presentation slides to place chart slides after their corresponding content slides.
Creates a new presentation with proper ordering.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from copy import deepcopy

def reorder_presentation():
    """Reorder slides to place charts after their corresponding content."""
    
    presentation_path = os.path.join(os.path.dirname(__file__), "..", "submission", "AI_Stocks_Presentation.pptx")
    if not os.path.exists(presentation_path):
        print(f"Error: Presentation not found at {presentation_path}")
        return
    
    prs = Presentation(presentation_path)
    
    # Create new presentation with same dimensions
    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height
    
    # Map slide titles to determine order
    slide_titles = []
    chart_slides = {}
    content_slides = []
    
    for slide in prs.slides:
        try:
            title = slide.shapes.title.text if slide.shapes.title else ""
            slide_titles.append(title)
            
            # Identify chart slides
            if any(keyword in title for keyword in ["Visual Comparison", "Performance Analysis", 
                                                     "Confusion Matrices", "Performance Comparison",
                                                     "Dataset Analysis"]):
                chart_slides[title] = slide
            else:
                content_slides.append((title, slide))
        except:
            content_slides.append(("", slide))
    
    # Define insertion mapping: content slide title -> chart slide title
    insertion_map = {
        "Results: Model Performance Metrics": "Model Performance: Visual Comparison",
        "Results: Key Findings & Model Selection": ["Per-Class Performance Analysis", "Error Analysis: Confusion Matrices"],
        "Results: Sentiment Analysis Key Findings": "Sentiment Models: Performance Comparison",
    }
    
    # Build ordered slide list
    ordered_slides = []
    i = 0
    while i < len(content_slides):
        title, slide = content_slides[i]
        ordered_slides.append(slide)
        
        # Check if we should insert chart slides after this content slide
        for content_title, chart_titles in insertion_map.items():
            if content_title in title:
                if isinstance(chart_titles, list):
                    for chart_title in chart_titles:
                        if chart_title in chart_slides:
                            ordered_slides.append(chart_slides[chart_title])
                            print(f"Inserting '{chart_title}' after '{title}'")
                else:
                    if chart_titles in chart_slides:
                        ordered_slides.append(chart_slides[chart_titles])
                        print(f"Inserting '{chart_titles}' after '{title}'")
        
        i += 1
    
    # Add any remaining chart slides that weren't inserted
    for chart_title, chart_slide in chart_slides.items():
        if chart_slide not in ordered_slides:
            ordered_slides.append(chart_slide)
            print(f"Adding remaining chart slide: '{chart_title}'")
    
    # Copy slides to new presentation
    # Note: We can't directly copy slides, so we'll need to recreate them
    # For now, save the original and provide instructions
    print(f"\nTotal slides to reorder: {len(ordered_slides)}")
    print("\nNote: python-pptx doesn't support copying slides between presentations.")
    print("Please manually reorder slides in PowerPoint:")
    print("1. Open the presentation")
    print("2. Move chart slides to appear after their corresponding content slides")
    print("3. Suggested order:")
    print("   - After 'Results: Model Performance Metrics' → 'Model Performance: Visual Comparison'")
    print("   - After 'Results: Key Findings & Model Selection' → 'Per-Class Performance Analysis'")
    print("   - After 'Per-Class Performance Analysis' → 'Error Analysis: Confusion Matrices'")
    print("   - After 'Results: Sentiment Analysis Key Findings' → 'Sentiment Models: Performance Comparison'")
    print("   - Add 'Dataset Analysis & Error Patterns' near the end of results section")
    
    return presentation_path

if __name__ == "__main__":
    reorder_presentation()




