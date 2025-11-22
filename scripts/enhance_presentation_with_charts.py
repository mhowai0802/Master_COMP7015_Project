#!/usr/bin/env python3
"""
Enhance the AI Stocks presentation with charts based on rubric requirements.
Adds comprehensive visualizations for model performance, comparisons, and error analysis.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import seaborn as sns
from io import BytesIO

# Set style for professional charts
try:
    plt.style.use('seaborn-v0_8-darkgrid')
except:
    try:
        plt.style.use('seaborn-darkgrid')
    except:
        plt.style.use('default')
sns.set_palette("husl")

def create_model_comparison_chart():
    """Create bar chart comparing MLP vs Transformer models."""
    fig, ax = plt.subplots(figsize=(10, 6))
    
    models = ['MLP', 'Transformer']
    metrics = ['Accuracy', 'Precision', 'Recall', 'F1-Score']
    
    mlp_scores = [40.93, 24.83, 32.50, 20.39]
    transformer_scores = [50.00, 34.75, 35.30, 29.35]
    
    x = np.arange(len(metrics))
    width = 0.35
    
    bars1 = ax.bar(x - width/2, mlp_scores, width, label='MLP', color='#3498db', alpha=0.8)
    bars2 = ax.bar(x + width/2, transformer_scores, width, label='Transformer', color='#e74c3c', alpha=0.8)
    
    ax.set_xlabel('Metrics', fontsize=12, fontweight='bold')
    ax.set_ylabel('Score (%)', fontsize=12, fontweight='bold')
    ax.set_title('Model Performance Comparison: MLP vs Transformer', fontsize=14, fontweight='bold', pad=20)
    ax.set_xticks(x)
    ax.set_xticklabels(metrics)
    ax.legend(fontsize=11)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    ax.set_ylim(0, 60)
    
    # Add value labels on bars
    for bars in [bars1, bars2]:
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height,
                   f'{height:.1f}%',
                   ha='center', va='bottom', fontsize=10, fontweight='bold')
    
    plt.tight_layout()
    return fig

def create_per_class_performance_chart():
    """Create per-class performance comparison."""
    fig, axes = plt.subplots(1, 2, figsize=(14, 6))
    
    classes = ['Down', 'Flat', 'Up']
    mlp_f1 = [58.0, 0.0, 4.0]
    transformer_f1 = [24.0, 0.0, 64.0]
    
    x = np.arange(len(classes))
    width = 0.35
    
    axes[0].bar(x - width/2, mlp_f1, width, label='MLP', color='#3498db', alpha=0.8)
    axes[0].bar(x + width/2, transformer_f1, width, label='Transformer', color='#e74c3c', alpha=0.8)
    axes[0].set_xlabel('Class', fontsize=11, fontweight='bold')
    axes[0].set_ylabel('F1-Score (%)', fontsize=11, fontweight='bold')
    axes[0].set_title('Per-Class F1-Score: MLP vs Transformer', fontsize=12, fontweight='bold')
    axes[0].set_xticks(x)
    axes[0].set_xticklabels(classes)
    axes[0].legend()
    axes[0].grid(axis='y', alpha=0.3, linestyle='--')
    axes[0].set_ylim(0, 70)
    
    # Add value labels
    for i, (m, t) in enumerate(zip(mlp_f1, transformer_f1)):
        if m > 0:
            axes[0].text(i - width/2, m, f'{m:.0f}%', ha='center', va='bottom', fontsize=9)
        if t > 0:
            axes[0].text(i + width/2, t, f'{t:.0f}%', ha='center', va='bottom', fontsize=9)
    
    # Precision and Recall comparison
    mlp_precision = [41.0, 0.0, 33.0]
    mlp_recall = [96.0, 0.0, 2.0]
    transformer_precision = [55.0, 0.0, 49.0]
    transformer_recall = [16.0, 0.0, 90.0]
    
    x2 = np.arange(len(classes))
    axes[1].bar(x2 - width*0.75, mlp_precision, width*0.5, label='MLP Precision', color='#3498db', alpha=0.6)
    axes[1].bar(x2 - width*0.25, mlp_recall, width*0.5, label='MLP Recall', color='#3498db', alpha=0.9)
    axes[1].bar(x2 + width*0.25, transformer_precision, width*0.5, label='Transformer Precision', color='#e74c3c', alpha=0.6)
    axes[1].bar(x2 + width*0.75, transformer_recall, width*0.5, label='Transformer Recall', color='#e74c3c', alpha=0.9)
    axes[1].set_xlabel('Class', fontsize=11, fontweight='bold')
    axes[1].set_ylabel('Score (%)', fontsize=11, fontweight='bold')
    axes[1].set_title('Per-Class Precision & Recall', fontsize=12, fontweight='bold')
    axes[1].set_xticks(x2)
    axes[1].set_xticklabels(classes)
    axes[1].legend(fontsize=9)
    axes[1].grid(axis='y', alpha=0.3, linestyle='--')
    axes[1].set_ylim(0, 100)
    
    plt.tight_layout()
    return fig

def create_confusion_matrix_mlp():
    """Create confusion matrix for MLP model."""
    fig, ax = plt.subplots(figsize=(8, 7))
    
    cm = np.array([[172, 0, 8],
                   [43, 0, 0],
                   [203, 0, 4]])
    
    sns.heatmap(cm, annot=True, fmt='d', cmap='Blues', ax=ax,
                xticklabels=['Down', 'Flat', 'Up'],
                yticklabels=['Down', 'Flat', 'Up'],
                cbar_kws={'label': 'Count'})
    
    ax.set_xlabel('Predicted', fontsize=12, fontweight='bold')
    ax.set_ylabel('Actual', fontsize=12, fontweight='bold')
    ax.set_title('MLP Model: Confusion Matrix', fontsize=14, fontweight='bold', pad=15)
    
    plt.tight_layout()
    return fig

def create_confusion_matrix_transformer():
    """Create confusion matrix for Transformer model."""
    fig, ax = plt.subplots(figsize=(8, 7))
    
    cm = np.array([[28, 0, 152],
                   [3, 0, 40],
                   [20, 0, 187]])
    
    sns.heatmap(cm, annot=True, fmt='d', cmap='Reds', ax=ax,
                xticklabels=['Down', 'Flat', 'Up'],
                yticklabels=['Down', 'Flat', 'Up'],
                cbar_kws={'label': 'Count'})
    
    ax.set_xlabel('Predicted', fontsize=12, fontweight='bold')
    ax.set_ylabel('Actual', fontsize=12, fontweight='bold')
    ax.set_title('Transformer Model: Confusion Matrix', fontsize=14, fontweight='bold', pad=15)
    
    plt.tight_layout()
    return fig

def create_sentiment_comparison_chart():
    """Create comparison chart for sentiment models (LSTM vs BERT)."""
    fig, axes = plt.subplots(1, 2, figsize=(14, 6))
    
    models = ['LSTM', 'BERT']
    
    # Accuracy comparison
    test_acc = [44.7, 50.0]
    val_acc = [44.7, 68.4]
    
    x = np.arange(len(models))
    width = 0.35
    
    axes[0].bar(x - width/2, test_acc, width, label='Test Accuracy', color='#9b59b6', alpha=0.8)
    axes[0].bar(x + width/2, val_acc, width, label='Validation Accuracy', color='#f39c12', alpha=0.8)
    axes[0].set_xlabel('Model', fontsize=11, fontweight='bold')
    axes[0].set_ylabel('Accuracy (%)', fontsize=11, fontweight='bold')
    axes[0].set_title('Sentiment Models: Accuracy Comparison', fontsize=12, fontweight='bold')
    axes[0].set_xticks(x)
    axes[0].set_xticklabels(models)
    axes[0].legend()
    axes[0].grid(axis='y', alpha=0.3, linestyle='--')
    axes[0].set_ylim(0, 80)
    
    # Add value labels
    for i, (test, val) in enumerate(zip(test_acc, val_acc)):
        axes[0].text(i - width/2, test, f'{test:.1f}%', ha='center', va='bottom', fontsize=10, fontweight='bold')
        axes[0].text(i + width/2, val, f'{val:.1f}%', ha='center', va='bottom', fontsize=10, fontweight='bold')
    
    # Macro metrics comparison
    lstm_metrics = [8.95, 20.00, 12.36]  # Precision, Recall, F1
    bert_metrics = [31.3, 31.0, 29.0]  # Approximate from results
    
    metrics = ['Precision', 'Recall', 'F1-Score']
    x2 = np.arange(len(metrics))
    
    axes[1].bar(x2 - width/2, lstm_metrics, width, label='LSTM', color='#9b59b6', alpha=0.8)
    axes[1].bar(x2 + width/2, bert_metrics, width, label='BERT', color='#f39c12', alpha=0.8)
    axes[1].set_xlabel('Metric', fontsize=11, fontweight='bold')
    axes[1].set_ylabel('Score (%)', fontsize=11, fontweight='bold')
    axes[1].set_title('Sentiment Models: Macro-Averaged Metrics', fontsize=12, fontweight='bold')
    axes[1].set_xticks(x2)
    axes[1].set_xticklabels(metrics)
    axes[1].legend()
    axes[1].grid(axis='y', alpha=0.3, linestyle='--')
    axes[1].set_ylim(0, 40)
    
    # Add value labels
    for bars in [axes[1].bar(x2 - width/2, lstm_metrics, width, color='#9b59b6', alpha=0.8),
                 axes[1].bar(x2 + width/2, bert_metrics, width, color='#f39c12', alpha=0.8)]:
        for bar in bars:
            height = bar.get_height()
            axes[1].text(bar.get_x() + bar.get_width()/2., height,
                        f'{height:.1f}%',
                        ha='center', va='bottom', fontsize=9, fontweight='bold')
    
    plt.tight_layout()
    return fig

def create_class_distribution_chart():
    """Create chart showing class distribution in test set."""
    fig, ax = plt.subplots(figsize=(8, 6))
    
    classes = ['Down', 'Flat', 'Up']
    counts = [180, 43, 207]
    colors = ['#e74c3c', '#95a5a6', '#2ecc71']
    
    bars = ax.bar(classes, counts, color=colors, alpha=0.8, edgecolor='black', linewidth=1.5)
    
    ax.set_xlabel('Class', fontsize=12, fontweight='bold')
    ax.set_ylabel('Number of Samples', fontsize=12, fontweight='bold')
    ax.set_title('Test Set Class Distribution', fontsize=14, fontweight='bold', pad=15)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    
    # Add value labels and percentages
    total = sum(counts)
    for bar, count in zip(bars, counts):
        height = bar.get_height()
        pct = (count / total) * 100
        ax.text(bar.get_x() + bar.get_width()/2., height,
               f'{count}\n({pct:.1f}%)',
               ha='center', va='bottom', fontsize=11, fontweight='bold')
    
    plt.tight_layout()
    return fig

def create_error_analysis_chart():
    """Create chart showing error patterns."""
    fig, ax = plt.subplots(figsize=(10, 6))
    
    error_types = ['Down→Up\n(MLP)', 'Down→Up\n(Trans)', 'Up→Down\n(MLP)', 'Up→Down\n(Trans)', 
                   'Flat→Down\n(MLP)', 'Flat→Up\n(Trans)']
    error_counts = [8, 152, 203, 20, 43, 40]
    colors = ['#e74c3c', '#c0392b', '#3498db', '#2980b9', '#95a5a6', '#7f8c8d']
    
    bars = ax.bar(range(len(error_types)), error_counts, color=colors, alpha=0.8, edgecolor='black', linewidth=1)
    
    ax.set_xlabel('Error Type', fontsize=12, fontweight='bold')
    ax.set_ylabel('Number of Misclassifications', fontsize=12, fontweight='bold')
    ax.set_title('Error Analysis: Common Misclassification Patterns', fontsize=14, fontweight='bold', pad=15)
    ax.set_xticks(range(len(error_types)))
    ax.set_xticklabels(error_types, rotation=45, ha='right')
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    
    # Add value labels
    for bar, count in zip(bars, error_counts):
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
               f'{count}',
               ha='center', va='bottom', fontsize=10, fontweight='bold')
    
    plt.tight_layout()
    return fig

def fig_to_image(fig):
    """Convert matplotlib figure to image bytes."""
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    return buf

def add_chart_to_slide(slide, fig, left, top, width, height):
    """Add a matplotlib figure to a PowerPoint slide."""
    img_stream = fig_to_image(fig)
    slide.shapes.add_picture(img_stream, left, top, width, height)
    plt.close(fig)

def enhance_presentation():
    """Enhance the existing presentation with charts."""
    
    # Load existing presentation
    presentation_path = os.path.join(os.path.dirname(__file__), "..", "submission", "AI_Stocks_Presentation.pptx")
    if not os.path.exists(presentation_path):
        print(f"Error: Presentation not found at {presentation_path}")
        return
    
    prs = Presentation(presentation_path)
    
    # Define colors
    title_color = RGBColor(41, 128, 185)
    accent_color = RGBColor(52, 73, 94)
    text_color = RGBColor(44, 62, 80)
    
    # Store slides to insert (index, slide)
    slides_to_insert = []
    
    # Find slide indices for insertion points
    performance_metrics_idx = None
    key_findings_idx = None
    sentiment_findings_idx = None
    
    for i, slide in enumerate(prs.slides):
        try:
            title_text = slide.shapes.title.text if slide.shapes.title else ""
            if "Model Performance Metrics" in title_text:
                performance_metrics_idx = i
            if "Key Findings & Model Selection" in title_text:
                key_findings_idx = i
            if "Sentiment Analysis Key Findings" in title_text:
                sentiment_findings_idx = i
        except:
            continue
    
    # Create new slides with charts
    # Use blank layout (index 6) or title-only (index 5)
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
    
    # 1. Model comparison chart after performance metrics
    if performance_metrics_idx is not None:
        slide = prs.slides.add_slide(blank_layout)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = "Model Performance: Visual Comparison"
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.color.rgb = title_color
        title_frame.paragraphs[0].font.bold = True
        
        fig = create_model_comparison_chart()
        add_chart_to_slide(slide, fig, Inches(0.5), Inches(1.3), Inches(9), Inches(5))
        slides_to_insert.append((performance_metrics_idx + 1, slide))
    
    # 2. Per-class performance after key findings
    if key_findings_idx is not None:
        slide = prs.slides.add_slide(blank_layout)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = "Per-Class Performance Analysis"
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.color.rgb = title_color
        title_frame.paragraphs[0].font.bold = True
        
        fig = create_per_class_performance_chart()
        add_chart_to_slide(slide, fig, Inches(0.3), Inches(1.3), Inches(9.4), Inches(5))
        slides_to_insert.append((key_findings_idx + 1, slide))
    
    # 3. Confusion matrices slide (add after per-class)
    slide = prs.slides.add_slide(blank_layout)
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Error Analysis: Confusion Matrices"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = title_color
    title_frame.paragraphs[0].font.bold = True
    
    # Add MLP confusion matrix
    fig1 = create_confusion_matrix_mlp()
    add_chart_to_slide(slide, fig1, Inches(0.5), Inches(1.3), Inches(4.2), Inches(4.5))
    
    # Add Transformer confusion matrix
    fig2 = create_confusion_matrix_transformer()
    add_chart_to_slide(slide, fig2, Inches(5.3), Inches(1.3), Inches(4.2), Inches(4.5))
    if key_findings_idx is not None:
        slides_to_insert.append((key_findings_idx + 2, slide))
    else:
        # Add at end if we can't find insertion point
        pass  # Already added to end
    
    # 4. Sentiment comparison after sentiment findings
    if sentiment_findings_idx is not None:
        slide = prs.slides.add_slide(blank_layout)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = "Sentiment Models: Performance Comparison"
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.color.rgb = title_color
        title_frame.paragraphs[0].font.bold = True
        
        fig = create_sentiment_comparison_chart()
        add_chart_to_slide(slide, fig, Inches(0.3), Inches(1.3), Inches(9.4), Inches(5))
        slides_to_insert.append((sentiment_findings_idx + 1, slide))
    
    # 5. Class distribution and error analysis
    slide = prs.slides.add_slide(blank_layout)
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Dataset Analysis & Error Patterns"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = title_color
    title_frame.paragraphs[0].font.bold = True
    
    # Add class distribution
    fig1 = create_class_distribution_chart()
    add_chart_to_slide(slide, fig1, Inches(0.5), Inches(1.3), Inches(4.2), Inches(4.5))
    
    # Add error analysis
    fig2 = create_error_analysis_chart()
    add_chart_to_slide(slide, fig2, Inches(5.3), Inches(1.3), Inches(4.2), Inches(4.5))
    # Add at end - will be after all other slides
    
    # Note: python-pptx doesn't support direct insertion, so slides are added at the end
    # User may need to manually reorder in PowerPoint, or we can create a new presentation
    # with proper ordering. For now, we'll add them at the end.
    
    # Save enhanced presentation
    output_path = presentation_path
    prs.save(output_path)
    print(f"Enhanced presentation saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print("\nNote: New chart slides have been added. You may want to reorder them in PowerPoint")
    print("to place them after their corresponding results slides for better flow.")
    
    return output_path

if __name__ == "__main__":
    try:
        enhance_presentation()
    except ImportError as e:
        print(f"Error: Missing required library. Please install: {e}")
        print("  pip install python-pptx matplotlib seaborn numpy")
    except Exception as e:
        print(f"Error enhancing presentation: {e}")
        import traceback
        traceback.print_exc()

